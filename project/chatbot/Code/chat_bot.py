import os
import csv
import concurrent.futures
from pathlib import Path
from pptx import Presentation
from pdfminer.high_level import extract_text
from langchain_community.embeddings import LlamaCppEmbeddings
from langchain_community.llms import GPT4All
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain.chains import ConversationalRetrievalChain
from langchain.memory import VectorStoreRetrieverMemory

CHUNK_SIZE = 500
CHUNK_OVERLAP = 200

def load_text(file_path):
    with open(file_path, 'r') as file:
        return {'page_content': file.read()}

def load_pdf(file_path):
    return {'page_content': extract_text(file_path)}

def load_csv(file_path):
    with open(file_path, newline='') as file:
        reader = csv.reader(file)
        return {'page_content': " ".join([" ".join(row) for row in reader])}

def load_pptx(file_path):
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return {'page_content': " ".join(text)}

def load_document(file_path):
    _, ext = os.path.splitext(file_path)
    loader_map = {
        '.pdf': load_pdf,
        '.txt': load_text,
        '.csv': load_csv,
        '.ppt': load_pptx,
        '.pptx': load_pptx,
    }
    loader = loader_map.get(ext.lower())
    if loader:
        return loader(file_path)
    return None

def process_document(file_path):
    document = load_document(file_path)
    if document:
        chunks = split_chunks([document])
        return chunks
    return []

def split_chunks(sources):
    chunks = []
    splitter = RecursiveCharacterTextSplitter(chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP)
    for source in sources:
        text = source.get('page_content', '')
        chunks.extend(splitter.split_text(text))
    return chunks

def generate_index(chunks, embeddings):
    texts = [chunk for chunk in chunks]
    metadatas = [{} for _ in chunks]
    return FAISS.from_texts(texts, embeddings, metadatas=metadatas)

class DocumentChatbot:
    def __init__(self):
        self.initialized = False
        self.embeddings = None
        self.llm = None
        self.index = None
        self.memory = None
        self.qa = None
        self.index_path = "/usr/src/app/Gpt4all/index"
        self.local_path = '/usr/src/app/Gpt4all/mistral-7b-openorca.gguf2.Q4_0.gguf'
        self.model_path = '/usr/src/app/Gpt4all/llama-2-7b.Q4_0.gguf'
        self.directory_path = "/usr/src/app/Pdf"

    def initialize(self):
        if self.initialized:
            return
        self.embeddings = LlamaCppEmbeddings(model_path=self.model_path)
        if os.path.exists(self.index_path):
            self.index = FAISS.load_local(self.index_path, self.embeddings)
        else:
            print("Index not found. Processing documents and generating index...")
            all_chunks = self.process_documents()
            self.index = generate_index(all_chunks, self.embeddings)
            self.index.save_local(self.index_path)

        if self.index is None:
            raise RuntimeError("Failed to load or generate index.")

        self.llm = GPT4All(model=self.local_path, verbose=True)
        self.qa = ConversationalRetrievalChain.from_llm(self.llm, retriever=self.index.as_retriever(), max_tokens_limit=400)
        self.initialized = True

    def process_documents(self):
        all_chunks = []
        with concurrent.futures.ThreadPoolExecutor() as executor:
            futures = [executor.submit(process_document, os.path.join(self.directory_path, filename))
                    for filename in os.listdir(self.directory_path) if os.path.isfile(os.path.join(self.directory_path, filename))]
            for future in concurrent.futures.as_completed(futures):
                all_chunks.extend(future.result())
        return all_chunks

    def process_query(self, query):
        self.initialize()  # Ensure the chatbot is initialized
        chat_history = []
        if not query:
            return "No query provided."
        result = self.qa({"question": query, "chat_history": chat_history})
        return result['answer']